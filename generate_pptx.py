from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Couleurs OCP ─────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1e, 0x3a, 0x5f)
MID_BLUE    = RGBColor(0x1e, 0x40, 0xaf)
LIGHT_BLUE  = RGBColor(0x93, 0xc5, 0xfd)
GREEN       = RGBColor(0x16, 0xa3, 0x4a)
ORANGE      = RGBColor(0xd9, 0x77, 0x06)
RED         = RGBColor(0xdc, 0x26, 0x26)
GRAY        = RGBColor(0x6b, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xf1, 0xf5, 0xf9)
WHITE       = RGBColor(0xff, 0xff, 0xff)
BLACK       = RGBColor(0x11, 0x18, 0x27)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout vide

# ── Helpers ───────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_textbox_multiline(slide, lines, l, t, w, h, size=14, color=BLACK,
                           bold_first=False, line_spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold_first and i == 0
    return txb

def header_bar(slide, title, subtitle=None):
    """Barre d'en-tête bleue foncée"""
    bar_h = 1.3 if subtitle else 1.1
    add_rect(slide, 0, 0, 13.33, bar_h, fill=DARK_BLUE)
    add_text(slide, title, 0.4, 0.1, 12, 0.7,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 12, 0.45,
                 size=14, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    # Ligne décorative orange
    add_rect(slide, 0, bar_h, 13.33, 0.05, fill=RGBColor(0xf5, 0x9e, 0x0b))

def footer(slide):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=DARK_BLUE)
    add_text(slide, "Maintenance Analytics  |  OCP Group — Site Daoui  |  Bureau de méthode",
             0.3, 7.22, 10, 0.25, size=9, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

def bullet(slide, items, l, t, w, size=14, color=BLACK, indent="  •  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(4.5))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if item.startswith("##"):
            run = p.add_run()
            run.text = item.replace("##", "").strip()
            run.font.size = Pt(size - 1)
            run.font.bold = True
            run.font.color.rgb = MID_BLUE
        elif item.startswith("#"):
            run = p.add_run()
            run.text = item.replace("#", "").strip()
            run.font.size = Pt(size + 1)
            run.font.bold = True
            run.font.color.rgb = DARK_BLUE
        else:
            run = p.add_run()
            run.text = indent + item
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return txb

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
# Fond dégradé simulé : deux rectangles
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(sl, 0, 4.5, 13.33, 3.0, fill=RGBColor(0x0f, 0x27, 0x4a))
# Bande accent
add_rect(sl, 0, 4.2, 13.33, 0.08, fill=RGBColor(0xf5, 0x9e, 0x0b))
# Logo texte
add_text(sl, "OCP GROUP — SITE DAOUI", 1, 0.6, 11, 0.5,
         size=13, color=LIGHT_BLUE, bold=False, align=PP_ALIGN.CENTER)
# Titre principal
add_text(sl, "Maintenance Analytics", 1, 1.3, 11.3, 1.5,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Sous-titre
add_text(sl, "Application web de pilotage de la maintenance industrielle", 1, 2.9, 11.3, 0.7,
         size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
# Séparateur
add_rect(sl, 3.5, 3.7, 6.3, 0.05, fill=RGBColor(0xf5, 0x9e, 0x0b))
# Infos bas
add_text(sl, "Bureau de méthode  |  SAP PM  |  Google Sheets  |  IA Chatbot", 1, 4.0, 11.3, 0.5,
         size=14, color=RGBColor(0xba, 0xd4, 0xf5), align=PP_ALIGN.CENTER)
add_text(sl, "2026", 1, 6.7, 11.3, 0.4,
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE & PROBLÉMATIQUE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Contexte & Problématique", "Pourquoi cette application a été développée ?")
footer(sl)

problems = [
    "Données SAP dispersées",
    "Les OTs, avis et PDR nécessitaient des exports manuels et des mises à jour périodiques",
    "Pas de visibilité en temps réel",
    "Impossible de connaître l'état des OTs LANC/CRPR sans générer un rapport SAP",
    "Suivi PDR 100% manuel",
    "Disponibilité des pièces gérée par emails et appels téléphoniques, sans traçabilité",
    "Demandes moteurs non structurées",
    "Demandes informelles : papier, WhatsApp, email — pas de suivi, pas d'historique",
    "Reporting chronophage",
    "Préparation des rapports hebdomadaires : 2 à 3 heures chaque semaine",
    "Coordination difficile entre services",
    "Mécanique, électrique, instrumentation, installation et bureau de méthode travaillaient en silos",
]

y = 1.45
for i in range(0, len(problems), 2):
    title_txt = problems[i]
    desc_txt  = problems[i+1] if i+1 < len(problems) else ""
    add_rect(sl, 0.4, y, 12.5, 0.08, fill=LIGHT_GRAY)
    add_text(sl, "⚠  " + title_txt, 0.5, y + 0.05, 12, 0.38,
             size=14, bold=True, color=RED)
    add_text(sl, desc_txt, 0.9, y + 0.42, 11.5, 0.38,
             size=12, color=GRAY)
    y += 0.9

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIFS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Objectifs de l'application", "4 axes stratégiques")
footer(sl)

objectives = [
    ("01", "Centralisation des données",
     "Connecter toutes les sources (SAP PM, Google Sheets) en un seul point d'accès visuel et interactif"),
    ("02", "Pilotage en temps réel",
     "KPIs actualisables à la demande : OTs LANC, CRPR, en retard, par poste, par installation"),
    ("03", "Digitalisation des processus",
     "Confirmation PDR en ligne • Demandes moteurs structurées • Workflow de validation et approbation"),
    ("04", "Communication automatisée",
     "Rappels automatiques chaque mercredi • Notifications push • Chatbot IA de pilotage instantané"),
]
colors = [MID_BLUE, GREEN, ORANGE, RGBColor(0x7c, 0x3a, 0xed)]

for i, (num, title, desc) in enumerate(objectives):
    col = 0.4 + (i % 2) * 6.45
    row = 1.45 + (i // 2) * 2.7
    add_rect(sl, col, row, 6.1, 2.4, fill=LIGHT_GRAY)
    add_rect(sl, col, row, 0.55, 2.4, fill=colors[i])
    add_text(sl, num, col + 0.05, row + 0.7, 0.5, 0.7,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, col + 0.7, row + 0.15, 5.2, 0.55,
             size=16, bold=True, color=colors[i])
    add_text(sl, desc, col + 0.7, row + 0.75, 5.2, 1.4,
             size=12, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE TECHNIQUE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Architecture Technique", "3 couches — données, logique, présentation")
footer(sl)

layers = [
    (MID_BLUE, "COUCHE DONNÉES — Google Sheets",
     ["Feuille 'Travaux hebdomadaire' : OTs SAP (statuts, postes, PDR, réalisation)",
      "Feuille 'Users' : gestion des utilisateurs et profils",
      "Feuilles 'Demandes' : moteurs électriques et interchangeables"]),
    (GREEN, "COUCHE LOGIQUE — Google Apps Script + Vercel",
     ["Google Apps Script : lecture/écriture Sheets, envoi d'emails automatiques",
      "Vercel (serverless) : proxy sécurisé entre l'application et les Google Sheets",
      "Aucun serveur à maintenir — infrastructure 100% cloud"]),
    (ORANGE, "COUCHE PRÉSENTATION — Application Web",
     ["Application web monopage (HTML / CSS / JavaScript)",
      "Interface responsive, accessible depuis n'importe quel navigateur",
      "Aucune installation requise — Chatbot IA intégré (Groq / LLaMA 3)"]),
]

for i, (color, title, items) in enumerate(layers):
    y = 1.45 + i * 1.8
    add_rect(sl, 0.3, y, 12.7, 1.65, fill=LIGHT_GRAY)
    add_rect(sl, 0.3, y, 3.2, 1.65, fill=color)
    add_text(sl, title, 0.35, y + 0.55, 3.1, 0.6,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, "•  " + item, 3.7, y + 0.1 + j * 0.48, 9.1, 0.45,
                 size=12, color=BLACK)

add_text(sl, "Navigateur  →  Vercel (proxy)  →  Google Apps Script  →  Google Sheets",
         1, 7.0, 11.3, 0.3, size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — PROFILS UTILISATEURS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Profils Utilisateurs", "7 profils — chaque service voit uniquement son périmètre")
footer(sl)

profiles = [
    ("Admin",                    "Accès complet",               "Validation demandes, paramètres globaux"),
    ("Responsable méthode",      "Lecture complète + CC emails", "Supervision globale de la maintenance"),
    ("Bureau de méthode",        "PDR spéciaux (keywords)",     "Réducteurs, pompes, moteurs"),
    ("Appro mécanique",          "PDR poste 421-MEC",           "Confirmation pièces mécaniques"),
    ("Appro électrique",         "PDR poste 423-ELEC",          "Confirmation pièces électriques"),
    ("Appro installation",       "PDR poste 421-INST",          "Confirmation pièces installation"),
    ("Appro Instrumentation",    "PDR poste 423-REG",           "Confirmation pièces instrumentation"),
    ("Interchangeable électrique","CC rappels + consultation",  "Suivi des pièces interchangeables"),
]

# En-tête tableau
header_y = 1.45
col_w = [3.5, 3.5, 5.1]
col_x = [0.3, 3.85, 7.4]
col_labels = ["Profil", "Accès", "Responsabilité"]
add_rect(sl, 0.3, header_y, 12.7, 0.42, fill=DARK_BLUE)
for ci, (cx, cw, cl) in enumerate(zip(col_x, col_w, col_labels)):
    add_text(sl, cl, cx + 0.1, header_y + 0.04, cw, 0.35,
             size=12, bold=True, color=WHITE)

for ri, (profil, acces, resp) in enumerate(profiles):
    ry = header_y + 0.45 + ri * 0.63
    bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
    add_rect(sl, 0.3, ry, 12.7, 0.6, fill=bg)
    add_text(sl, profil, col_x[0] + 0.1, ry + 0.1, col_w[0], 0.4,
             size=11, bold=True, color=DARK_BLUE)
    add_text(sl, acces,  col_x[1] + 0.1, ry + 0.1, col_w[1], 0.4,
             size=11, color=BLACK)
    add_text(sl, resp,   col_x[2] + 0.1, ry + 0.1, col_w[2], 0.4,
             size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — TABLEAU DE BORD
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Tableau de Bord", "KPIs en temps réel — vue d'ensemble instantanée")
footer(sl)

kpis = [
    (MID_BLUE,  "OTs LANC",   "Lancés / en cours d'exécution"),
    (ORANGE,    "OTs CRPR",   "En cours de préparation"),
    (RED,       "OTs Retard", "Non clôturés hors délai"),
    (GREEN,     "OTs CONF",   "Confirmés / clôturés"),
    (RGBColor(0x7c,0x3a,0xed), "Avis AOUV", "En attente de création OT"),
    (RGBColor(0x0e,0x7a,0x90), "Avis AENC", "En cours de traitement"),
]

for i, (color, label, desc) in enumerate(kpis):
    col = 0.35 + (i % 3) * 4.3
    row = 1.45 + (i // 3) * 2.2
    add_rect(sl, col, row, 4.0, 1.9, fill=LIGHT_GRAY)
    add_rect(sl, col, row, 4.0, 0.5, fill=color)
    add_text(sl, label, col + 0.1, row + 0.06, 3.8, 0.42,
             size=15, bold=True, color=WHITE)
    add_text(sl, "███", col + 0.2, row + 0.65, 1.2, 0.5,
             size=28, bold=True, color=color)
    add_text(sl, desc, col + 0.15, row + 1.35, 3.7, 0.4,
             size=11, color=GRAY)

add_text(sl, "🔄  Bouton d'actualisation en haut à droite — rafraîchit toutes les données sans quitter la page en cours",
         0.3, 6.85, 12.7, 0.35, size=11, color=MID_BLUE, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — MODULE OTs SAP
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Module Rapport — OTs SAP", "Consultation complète des ordres de travail de la semaine")
footer(sl)

add_rect(sl, 0.3, 1.45, 5.9, 5.4, fill=LIGHT_GRAY)
add_rect(sl, 6.5, 1.45, 6.5, 5.4, fill=LIGHT_GRAY)

add_text(sl, "Filtres disponibles", 0.5, 1.55, 5.5, 0.45, size=15, bold=True, color=DARK_BLUE)
filtres = ["Par poste de travail (421-MEC, 423-ELEC...)",
           "Par statut système (LANC, CRPR, CONF, CLOT...)",
           "Par statut utilisateur (ATPL, CRPR...)",
           "Par installation",
           "Recherche libre par mot-clé"]
for i, f in enumerate(filtres):
    add_text(sl, "✓  " + f, 0.5, 2.1 + i * 0.55, 5.5, 0.45, size=12, color=BLACK)

add_text(sl, "Informations par OT", 6.7, 1.55, 6.0, 0.45, size=15, bold=True, color=DARK_BLUE)
infos = ["N° d'ordre, objet technique, description",
         "Poste de travail et installation",
         "Statut SAP (système + utilisateur)",
         "PDR associée et statut de confirmation",
         "Réalisation (Fait / Non fait)",
         "Export et impression pour rapports hebdomadaires"]
for i, f in enumerate(infos):
    add_text(sl, "→  " + f, 6.7, 2.1 + i * 0.55, 6.0, 0.45, size=12, color=BLACK)

add_rect(sl, 0.3, 6.5, 12.7, 0.5, fill=RGBColor(0xe0, 0xf2, 0xfe))
add_text(sl, "💡  Bénéfice : accès immédiat à tous les OTs SAP sans connexion à SAP ni export manuel",
         0.5, 6.56, 12.3, 0.35, size=12, color=MID_BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — MODULE AVIS SAP
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Module Avis SAP", "Suivi des avis en attente de traitement")
footer(sl)

add_rect(sl, 0.3, 1.45, 12.7, 0.55, fill=DARK_BLUE)
for ci, label in enumerate(["Statut", "Signification", "Action requise"]):
    add_text(sl, label, 0.5 + ci * 4.2, 1.5, 4.0, 0.4,
             size=13, bold=True, color=WHITE)

rows_avis = [
    ("AOUV", "Avis ouvert", "Créer un OT dans SAP"),
    ("AENC", "Avis en cours de traitement", "Suivre l'avancement"),
    ("CLOT", "Avis clôturé", "Historique — aucune action"),
]
for ri, (statut, sig, action) in enumerate(rows_avis):
    ry = 2.05 + ri * 0.7
    bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
    add_rect(sl, 0.3, ry, 12.7, 0.65, fill=bg)
    color = RED if statut == "AOUV" else (ORANGE if statut == "AENC" else GREEN)
    add_rect(sl, 0.3, ry, 1.5, 0.65, fill=color)
    add_text(sl, statut, 0.35, ry + 0.12, 1.4, 0.4,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, sig,    1.9,  ry + 0.12, 3.9, 0.4, size=12, color=BLACK)
    add_text(sl, action, 6.1,  ry + 0.12, 6.6, 0.4, size=12, color=GRAY)

add_text(sl, "Indicateurs affichés", 0.5, 3.6, 12.0, 0.45, size=15, bold=True, color=DARK_BLUE)
inds = ["Nombre total d'avis AOUV + AENC (en attente OT)",
        "Répartition par poste de travail et par installation",
        "Liste des avis les plus anciens — priorités de traitement"]
for i, ind in enumerate(inds):
    add_text(sl, "•  " + ind, 0.5, 4.15 + i * 0.55, 12, 0.45, size=13, color=BLACK)

add_rect(sl, 0.3, 6.5, 12.7, 0.5, fill=RGBColor(0xe0, 0xf2, 0xfe))
add_text(sl, "💡  Plus aucun avis oublié dans SAP — visibilité temps réel sur tous les avis du site",
         0.5, 6.56, 12.3, 0.35, size=12, color=MID_BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — MODULE PDR
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Module PDR — Pièces de Rechange", "Le cœur opérationnel de l'application")
footer(sl)

add_rect(sl, 0.3, 1.45, 12.7, 1.0, fill=RGBColor(0xff, 0xf7, 0xed))
add_text(sl, "⚠  Avant :", 0.5, 1.5, 2, 0.4, size=13, bold=True, color=ORANGE)
add_text(sl, "La disponibilité des PDR se gérait par téléphone et email. Les OTs étaient bloqués faute de confirmation à temps. Pas de traçabilité.",
         2.5, 1.5, 10.2, 0.8, size=12, color=BLACK)

add_text(sl, "Affectation intelligente par poste", 0.5, 2.65, 12, 0.45, size=15, bold=True, color=DARK_BLUE)

postes = [
    (MID_BLUE,  "421-MEC",  "Appro mécanique"),
    (GREEN,     "423-ELEC", "Appro électrique"),
    (ORANGE,    "421-INST", "Appro installation"),
    (RGBColor(0x7c,0x3a,0xed), "423-REG", "Appro Instrumentation"),
    (RED,       "Réducteur / Pompe / Moteur", "Bureau de méthode"),
]
for i, (color, poste, service) in enumerate(postes):
    col = 0.3 + (i % 3) * 4.3 if i < 3 else 0.3 + (i - 3) * 6.5
    row = 3.2 if i < 3 else 4.85
    w = 3.9 if i < 3 else 6.0
    add_rect(sl, col, row, w, 1.35, fill=LIGHT_GRAY)
    add_rect(sl, col, row, w, 0.38, fill=color)
    add_text(sl, poste,   col + 0.1, row + 0.04, w - 0.2, 0.3,
             size=13, bold=True, color=WHITE)
    add_text(sl, "→ " + service, col + 0.1, row + 0.48, w - 0.2, 0.75,
             size=12, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — PDR STATUTS DE CONFIRMATION
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "PDR — Statuts de Confirmation", "4 statuts pour une traçabilité complète")
footer(sl)

statuts = [
    (GREEN,  "✅  En stock",       "La pièce est disponible en stock.\nConfirmation immédiate — OT peut démarrer."),
    (RED,    "❌  Non disponible", "La pièce est manquante.\nDélai de livraison estimé renseigné."),
    (ORANGE, "💬  Observation",    "Information partielle — complément nécessaire.\nSauvegardable SANS choisir un statut."),
    (GRAY,   "⏳  En attente",    "Aucune réponse encore.\nRappel automatique envoyé chaque mercredi."),
]

for i, (color, label, desc) in enumerate(statuts):
    col = 0.3 + (i % 2) * 6.5
    row = 1.45 + (i // 2) * 2.6
    add_rect(sl, col, row, 6.2, 2.3, fill=LIGHT_GRAY)
    add_rect(sl, col, row, 6.2, 0.55, fill=color)
    add_text(sl, label, col + 0.15, row + 0.07, 5.9, 0.42,
             size=16, bold=True, color=WHITE)
    for li, line in enumerate(desc.split("\n")):
        add_text(sl, line, col + 0.2, row + 0.75 + li * 0.55, 5.8, 0.45,
                 size=13, color=BLACK)

add_rect(sl, 0.3, 6.9, 12.7, 0.18, fill=RGBColor(0xf5, 0x9e, 0x0b))

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — PDR WORKFLOW
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "PDR — Workflow de Confirmation", "De la détection à la confirmation en quelques clics")
footer(sl)

steps = [
    (MID_BLUE,  "1", "Export SAP",         "Les données SAP PM sont exportées dans\nGoogle Sheets chaque semaine"),
    (GREEN,     "2", "Détection auto",      "L'application identifie tous les OTs\navec PDR renseignée"),
    (ORANGE,    "3", "Affichage service",   "Chaque service voit uniquement\nses PDR à confirmer"),
    (RED,       "4", "Confirmation",        "Le service appro confirme : En stock,\nNon disponible ou Observation"),
    (RGBColor(0x7c,0x3a,0xed), "5", "Sauvegarde", "Donnée enregistrée dans Google Sheets\n(colonnes T et U)"),
    (DARK_BLUE, "6", "Visibilité immédiate","Bureau de méthode voit la confirmation\nen temps réel"),
]

for i, (color, num, title, desc) in enumerate(steps):
    col = 0.3 + (i % 3) * 4.3
    row = 1.45 + (i // 3) * 2.6
    add_rect(sl, col, row, 4.0, 2.3, fill=LIGHT_GRAY)
    add_rect(sl, col, row, 0.6, 2.3, fill=color)
    add_text(sl, num, col + 0.05, row + 0.75, 0.52, 0.6,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, col + 0.75, row + 0.12, 3.1, 0.5,
             size=14, bold=True, color=color)
    for li, line in enumerate(desc.split("\n")):
        add_text(sl, line, col + 0.75, row + 0.7 + li * 0.48, 3.1, 0.45,
                 size=11, color=BLACK)

    if i % 3 < 2:
        add_text(sl, "→", col + 4.1, row + 0.9, 0.3, 0.5,
                 size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — RAPPELS AUTOMATIQUES PAR EMAIL
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Rappels Automatiques par Email", "Chaque mercredi à 08h00 — sans intervention humaine")
footer(sl)

add_rect(sl, 0.3, 1.45, 12.7, 1.1, fill=RGBColor(0xe0, 0xf7, 0xea))
add_text(sl, "🕗  Déclenchement automatique : chaque mercredi à 08h00 via Google Apps Script",
         0.5, 1.55, 12.3, 0.4, size=14, bold=True, color=GREEN)
add_text(sl, "Le script scanne tous les OTs actifs (LANC/CRPR) avec PDR sans confirmation et envoie les rappels.",
         0.5, 1.95, 12.3, 0.45, size=12, color=BLACK)

add_text(sl, "Destinataires par poste", 0.5, 2.75, 6.0, 0.45, size=15, bold=True, color=DARK_BLUE)
add_text(sl, "CC systématique", 6.8, 2.75, 6.0, 0.45, size=15, bold=True, color=DARK_BLUE)

dest_rows = [
    ("421-MEC",  "Appro mécanique"),
    ("423-ELEC", "Appro électrique"),
    ("421-INST", "Appro installation"),
    ("423-REG",  "Appro Instrumentation"),
]
for i, (poste, service) in enumerate(dest_rows):
    ry = 3.3 + i * 0.62
    add_rect(sl, 0.3, ry, 6.2, 0.55, fill=LIGHT_GRAY if i%2==0 else WHITE)
    add_text(sl, poste,   0.4, ry + 0.08, 2.0, 0.38, size=12, bold=True, color=MID_BLUE)
    add_text(sl, "→ " + service, 2.4, ry + 0.08, 3.8, 0.38, size=12, color=BLACK)

cc_items = ["Responsable méthode", "Interchangeable électrique"]
for i, cc in enumerate(cc_items):
    add_rect(sl, 6.8, 3.3 + i * 0.7, 6.0, 0.6, fill=LIGHT_GRAY)
    add_text(sl, "📋  " + cc, 6.95, 3.38 + i * 0.7, 5.6, 0.42, size=13, color=DARK_BLUE)

add_text(sl, "Contenu de l'email :", 0.5, 5.75, 4.0, 0.4, size=13, bold=True, color=DARK_BLUE)
add_text(sl, "Tableau récapitulatif : N° OT  •  Objet technique  •  Description  •  PDR demandée",
         0.5, 6.2, 12.5, 0.38, size=12, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — DEMANDES MOTEURS ÉLECTRIQUES
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Module Demandes Moteurs Électriques", "Digitalisation des demandes de remplacement")
footer(sl)

add_rect(sl, 0.3, 1.45, 5.9, 5.35, fill=LIGHT_GRAY)
add_rect(sl, 6.5, 1.45, 6.5, 5.35, fill=LIGHT_GRAY)

add_text(sl, "Formulaire de demande", 0.5, 1.55, 5.5, 0.45, size=15, bold=True, color=DARK_BLUE)
form_fields = [
    "Type : Pose (nouveau) ou Réparation",
    "Installation et objet technique",
    "Puissance (kW), tension (V), vitesse (tr/min)",
    "Description de l'anomalie",
    "Matricule et nom du demandeur",
]
for i, f in enumerate(form_fields):
    add_text(sl, "📝  " + f, 0.5, 2.12 + i * 0.6, 5.6, 0.5, size=12, color=BLACK)

add_text(sl, "Workflow de validation", 6.7, 1.55, 6.0, 0.45, size=15, bold=True, color=DARK_BLUE)

add_rect(sl, 6.7, 2.12, 6.0, 0.95, fill=RGBColor(0xe0, 0xf2, 0xfe))
add_text(sl, "Pose (nouveau moteur)", 6.9, 2.18, 5.6, 0.35, size=13, bold=True, color=MID_BLUE)
add_text(sl, "→ Approbation Admin requise\n→ Email de notification envoyé",
         6.9, 2.52, 5.6, 0.45, size=11, color=BLACK)

add_rect(sl, 6.7, 3.22, 6.0, 0.95, fill=RGBColor(0xe0, 0xf7, 0xea))
add_text(sl, "Réparation", 6.9, 3.28, 5.6, 0.35, size=13, bold=True, color=GREEN)
add_text(sl, "→ Approbation automatique\n→ Urgence opérationnelle prioritaire",
         6.9, 3.62, 5.6, 0.45, size=11, color=BLACK)

add_text(sl, "Suivi en temps réel", 6.7, 4.35, 6.0, 0.4, size=14, bold=True, color=DARK_BLUE)
suivi = ["Statut : En attente / Approuvé / Refusé",
         "Historique complet avec dates",
         "Justification admin en cas de refus"]
for i, s in enumerate(suivi):
    add_text(sl, "•  " + s, 6.7, 4.82 + i * 0.5, 6.0, 0.42, size=12, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — CHATBOT IA
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Chatbot IA Intégré", "Un assistant intelligent disponible 24h/24 — Powered by Groq / LLaMA 3")
footer(sl)

add_rect(sl, 0.3, 1.45, 5.9, 5.35, fill=LIGHT_GRAY)
add_rect(sl, 6.5, 1.45, 6.5, 5.35, fill=LIGHT_GRAY)

add_text(sl, "Ce que le chatbot connaît", 0.5, 1.55, 5.5, 0.45, size=15, bold=True, color=DARK_BLUE)
knows = ["Tous les KPIs du tableau de bord",
         "OTs LANC et CRPR — répartition par poste",
         "Statut des PDR (en attente / confirmées)",
         "Données des avis SAP (AOUV, AENC)",
         "Planning de la semaine en cours",
         "Contexte mis à jour à chaque actualisation"]
for i, k in enumerate(knows):
    add_text(sl, "✓  " + k, 0.5, 2.1 + i * 0.55, 5.6, 0.45, size=12, color=BLACK)

add_text(sl, "Exemples de questions", 6.7, 1.55, 6.0, 0.45, size=15, bold=True, color=DARK_BLUE)
questions = [
    '"Combien d\'OTs sont en attente de PDR ?"',
    '"Répartition des OTs LANC par corps de métier ?"',
    '"Quels sont les OTs CRPR du poste 421-MEC ?"',
    '"Combien d\'avis sont en attente de traitement ?"',
]
for i, q in enumerate(questions):
    add_rect(sl, 6.7, 2.12 + i * 1.1, 6.0, 0.9, fill=RGBColor(0xee, 0xf2, 0xff))
    add_text(sl, q, 6.85, 2.2 + i * 1.1, 5.7, 0.7, size=12, color=MID_BLUE, italic=True)

add_text(sl, "🤖  Icône robot animé flottante — accessible depuis toutes les pages de l'application",
         0.3, 6.85, 12.7, 0.3, size=11, color=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — ARRÊTS PLANIFIÉS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Module Arrêts Planifiés", "Suivi et pilotage des arrêts de production")
footer(sl)

features = [
    (MID_BLUE,  "Calendrier des arrêts",
     "Vue mensuelle par installation\nDurée, type, équipement concerné\nStatut : planifié / en cours / terminé / reporté"),
    (ORANGE,    "Marquage Reporté",
     "Arrêt marqué 'Reporté' directement depuis l'app\nMise à jour automatique dans Google Sheets\nTraçabilité complète avec date de report"),
    (GREEN,     "Indicateurs de pilotage",
     "Nombre d'arrêts planifiés sur la période\nTaux de réalisation vs planification\nArrêts reportés — suivi des glissements"),
]

for i, (color, title, desc) in enumerate(features):
    col = 0.3 + i * 4.3
    add_rect(sl, col, 1.45, 4.1, 4.8, fill=LIGHT_GRAY)
    add_rect(sl, col, 1.45, 4.1, 0.55, fill=color)
    add_text(sl, title, col + 0.1, 1.51, 3.9, 0.42,
             size=14, bold=True, color=WHITE)
    for li, line in enumerate(desc.split("\n")):
        add_text(sl, "•  " + line, col + 0.15, 2.15 + li * 0.62, 3.8, 0.5,
                 size=12, color=BLACK)

add_rect(sl, 0.3, 6.5, 12.7, 0.5, fill=RGBColor(0xe0, 0xf2, 0xfe))
add_text(sl, "💡  Les équipes terrain peuvent signaler un report sans passer par une réunion ou un email — traçabilité instantanée",
         0.5, 6.56, 12.3, 0.35, size=11, color=MID_BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — SUIVI DES RÉALISATIONS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Suivi des Réalisations", "Clôture des OTs en temps réel — sans passer par SAP")
footer(sl)

add_rect(sl, 0.3, 1.45, 12.7, 1.4, fill=LIGHT_GRAY)
add_text(sl, "Comment ça fonctionne ?", 0.5, 1.52, 12.3, 0.42, size=15, bold=True, color=DARK_BLUE)
add_text(sl, "Le technicien ou chef d'équipe ouvre l'OT dans l'application → sélectionne Fait ou Non Fait → la valeur est enregistrée dans Google Sheets (colonne O) → visible immédiatement par le bureau de méthode.",
         0.5, 1.98, 12.3, 0.7, size=12, color=BLACK)

stats_items = [
    (GREEN, "✅  FAIT",     "OT réalisé dans les délais\nMarqué et tracé avec horodatage"),
    (RED,   "❌  NON FAIT", "OT non réalisé\nAnalyse des causes facilitée"),
]
for i, (color, label, desc) in enumerate(stats_items):
    col = 0.3 + i * 6.5
    add_rect(sl, col, 3.1, 6.2, 2.1, fill=LIGHT_GRAY)
    add_rect(sl, col, 3.1, 6.2, 0.55, fill=color)
    add_text(sl, label, col + 0.1, 3.16, 6.0, 0.42, size=16, bold=True, color=WHITE)
    for li, line in enumerate(desc.split("\n")):
        add_text(sl, line, col + 0.2, 3.78 + li * 0.48, 5.8, 0.42, size=12, color=BLACK)

add_text(sl, "Bénéfices", 0.5, 5.45, 12.3, 0.42, size=15, bold=True, color=DARK_BLUE)
benefits = ["Taux de réalisation hebdomadaire calculé automatiquement",
            "Identification rapide des OTs non réalisés — plan d'action immédiat",
            "Rapport de réalisation généré sans saisie manuelle supplémentaire"]
for i, b in enumerate(benefits):
    add_text(sl, "→  " + b, 0.5, 5.98 + i * 0.47, 12.3, 0.4, size=12, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — NOTIFICATIONS PUSH
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Notifications Push", "Rester informé en temps réel — même sans avoir l'application ouverte")
footer(sl)

triggers_notif = [
    (MID_BLUE,  "Nouvelle demande moteur soumise",  "Admin notifié immédiatement"),
    (GREEN,     "Demande approuvée",                 "Demandeur notifié avec confirmation"),
    (RED,       "Demande refusée",                   "Demandeur notifié avec justification"),
    (ORANGE,    "PDR confirmée",                     "Bureau de méthode notifié"),
]

add_text(sl, "Déclencheurs de notification", 0.5, 1.45, 8.0, 0.45, size=15, bold=True, color=DARK_BLUE)
for i, (color, trigger, action) in enumerate(triggers_notif):
    ry = 2.0 + i * 0.85
    add_rect(sl, 0.3, ry, 12.7, 0.75, fill=LIGHT_GRAY if i%2==0 else WHITE)
    add_rect(sl, 0.3, ry, 0.2, 0.75, fill=color)
    add_text(sl, "🔔  " + trigger, 0.6, ry + 0.15, 7.0, 0.42, size=13, bold=True, color=color)
    add_text(sl, action, 7.8, ry + 0.15, 5.0, 0.42, size=12, color=GRAY)

add_text(sl, "Caractéristiques techniques", 0.5, 5.55, 12.3, 0.42, size=15, bold=True, color=DARK_BLUE)
techs = ["Notifications push navigateur — fonctionne même si l'onglet est fermé",
         "Aucune application mobile nécessaire — compatible PC et mobile",
         "Abonnement par profil utilisateur — chaque service reçoit ses notifications"]
for i, t in enumerate(techs):
    add_text(sl, "✓  " + t, 0.5, 6.05 + i * 0.42, 12.3, 0.38, size=12, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — SÉCURITÉ
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Sécurité & Accès", "Une application sécurisée — données protégées")
footer(sl)

security_items = [
    (MID_BLUE, "Authentification",
     ["Accès par profil utilisateur géré dans la feuille Users",
      "Chaque utilisateur voit uniquement les données de son périmètre",
      "Pas d'accès croisé entre services"]),
    (GREEN, "Protection des données",
     ["Feuille Users : accès restreint (non public)",
      "Dépôt de code source : GitHub privé",
      "Clés API gérées côté serveur — non exposées dans le code client"]),
    (ORANGE, "Infrastructure sécurisée",
     ["Hébergement Vercel : HTTPS obligatoire sur tous les échanges",
      "Google Apps Script : exécution avec compte autorisé uniquement",
      "Aucune donnée stockée hors de l'environnement OCP Google Workspace"]),
]

for i, (color, title, items) in enumerate(security_items):
    col = 0.3 + i * 4.3
    add_rect(sl, col, 1.45, 4.1, 5.2, fill=LIGHT_GRAY)
    add_rect(sl, col, 1.45, 4.1, 0.55, fill=color)
    add_text(sl, "🔒  " + title, col + 0.1, 1.51, 3.9, 0.42,
             size=13, bold=True, color=WHITE)
    for li, item in enumerate(items):
        add_text(sl, "•  " + item, col + 0.15, 2.15 + li * 0.85, 3.8, 0.75,
                 size=11, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — BÉNÉFICES & GAINS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Bénéfices & Gains Mesurables", "Ce que Maintenance Analytics change concrètement")
footer(sl)

add_text(sl, "Gains de temps", 0.4, 1.45, 12.5, 0.45, size=15, bold=True, color=DARK_BLUE)

# Tableau des gains
col_headers = ["Tâche", "Avant", "Après"]
col_x2 = [0.3, 5.8, 9.8]
col_w2 = [5.4, 3.8, 3.3]
add_rect(sl, 0.3, 1.95, 12.7, 0.42, fill=DARK_BLUE)
for ci, (cx, cw, cl) in enumerate(zip(col_x2, col_w2, col_headers)):
    add_text(sl, cl, cx + 0.1, 2.0, cw, 0.32,
             size=12, bold=True, color=WHITE)

gain_rows = [
    ("Rapport hebdomadaire OTs",  "2-3 heures",        "5 minutes (actualisation)"),
    ("Suivi PDR par poste",       "30 min emails/appels", "Temps réel sur l'écran"),
    ("Demande de moteur",         "Formulaire papier", "Formulaire numérique direct"),
    ("Rappel confirmation PDR",   "Manuel chaque semaine", "Automatique chaque mercredi"),
    ("Réponse pilotage",          "Consultation SAP",  "Chatbot instantané"),
]
for ri, (task, before, after) in enumerate(gain_rows):
    ry = 2.4 + ri * 0.6
    bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
    add_rect(sl, 0.3, ry, 12.7, 0.55, fill=bg)
    add_text(sl, task,   col_x2[0] + 0.1, ry + 0.08, col_w2[0], 0.38, size=11, color=BLACK)
    add_text(sl, "⏱  " + before, col_x2[1] + 0.1, ry + 0.08, col_w2[1], 0.38, size=11, color=RED)
    add_text(sl, "✅  " + after,  col_x2[2] + 0.1, ry + 0.08, col_w2[2], 0.38, size=11, color=GREEN, bold=True)

add_text(sl, "Gains qualitatifs", 0.4, 5.6, 12.5, 0.38, size=14, bold=True, color=DARK_BLUE)
qual = ["Traçabilité complète de toutes les actions (PDR, demandes, réalisations)",
        "Coordination améliorée entre services • Réactivité augmentée grâce aux notifications",
        "Zéro email informel pour demandes et confirmations PDR"]
for i, q in enumerate(qual):
    add_text(sl, "→  " + q, 0.4, 6.05 + i * 0.38, 12.5, 0.33, size=11, color=BLACK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — CONCLUSION & PERSPECTIVES
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Conclusion & Perspectives", "Une base solide — des évolutions continues")
footer(sl)

add_rect(sl, 0.3, 1.45, 6.1, 5.35, fill=LIGHT_GRAY)
add_rect(sl, 6.7, 1.45, 6.3, 5.35, fill=LIGHT_GRAY)

add_rect(sl, 0.3, 1.45, 6.1, 0.5, fill=GREEN)
add_text(sl, "✅  Opérationnel aujourd'hui", 0.45, 1.51, 5.8, 0.38,
         size=14, bold=True, color=WHITE)

done = ["Tableau de bord temps réel (OTs, Avis, KPIs)",
        "Suivi PDR avec confirmation multi-service",
        "Rappels automatiques email (4 services + CC)",
        "Demandes moteurs avec workflow de validation",
        "Chatbot IA de pilotage (Groq / LLaMA 3)",
        "Notifications push navigateur",
        "Suivi des réalisations hebdomadaires",
        "Calendrier des arrêts planifiés"]
for i, d in enumerate(done):
    add_text(sl, "✓  " + d, 0.45, 2.08 + i * 0.56, 5.8, 0.48, size=11, color=BLACK)

add_rect(sl, 6.7, 1.45, 6.3, 0.5, fill=MID_BLUE)
add_text(sl, "🚀  Évolutions envisageables", 6.85, 1.51, 6.0, 0.38,
         size=14, bold=True, color=WHITE)

future = ["Tableaux analytiques avancés (tendances multi-semaines)",
          "Application mobile native",
          "Connexion directe SAP (sans export manuel)",
          "Génération automatique de rapports PDF",
          "Escalade auto pour PDR non confirmées après X jours",
          "Module gestion de stock pièces de rechange",
          "Tableau de bord direction avec synthèse mensuelle"]
for i, f in enumerate(future):
    add_text(sl, "◦  " + f, 6.85, 2.08 + i * 0.56, 6.0, 0.48, size=11, color=BLACK)

add_rect(sl, 0.3, 7.0, 12.7, 0.2, fill=RGBColor(0xf5, 0x9e, 0x0b))

# ── Sauvegarde ───────────────────────────────────────────────
output_path = '/Users/mounaim/Desktop/mainAna/Presentation_MaintenanceAnalytics.pptx'
prs.save(output_path)
print("PPTX généré :", output_path)
